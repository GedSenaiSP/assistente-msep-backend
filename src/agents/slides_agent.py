"""
Agente LangGraph para geração de slides PowerPoint.

Gera apresentações educacionais baseadas em Situações de Aprendizagem,
utilizando o template SENAI fornecido.
"""
import os
import json
import logging
from io import BytesIO
from typing import Dict, List, Any, TypedDict
from pathlib import Path

from langgraph.graph import StateGraph, END
from langchain_openai import AzureChatOpenAI

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Configuração de logs
logger = logging.getLogger(__name__)

# Configurações de templates disponíveis
TEMPLATES = {
    "dn": {
        "name": "Departamento Nacional",
        "path": Path(__file__).parent.parent / "template" / "template-slides-dn.pptx",
        # Template com 10 layouts
        "layout_capa": 0,           # Capa (visual, sem placeholders)
        "layout_titulo": 1,         # Slide de Título (Contracapa do início)
        "layout_capa_secao": 2,     # Capa da Seção (Divisor)
        "layout_conteudo": 7,       # Título e Conteúdo
        "layout_fim": 3,            # Encerramento Visual
        "placeholder_title": [0],       # Índice para título em slides de conteúdo
        "placeholder_content": 10,      # Índice para conteúdo
        "placeholder_capa": 0,          # Índice para texto da capa (no layout_capa_secao) - DN layout 2 tem texto no 0
    },
    "sp": {
        "name": "São Paulo",
        "path": Path(__file__).parent.parent / "template" / "template-slides-sp.pptx",
        # Template com 7 layouts
        "layout_capa": 0,           # Capa (visual)
        "layout_titulo": 1,         # Slide de Título (Contracapa)
        "layout_capa_secao": 2,     # Capa da Seção
        "layout_conteudo": 5,       # Título e Conteúdo
        "layout_fim": 6,            # Encerramento
        "placeholder_title": [10],      # Índice para título (BODY)
        "placeholder_content": 11,      # Índice para conteúdo (OBJECT)
        "placeholder_capa": 14,         # Índice para texto da capa/seção
        "placeholder_section_title": 14, # Placeholder específico para título de seção
    }
}

# Template padrão
DEFAULT_TEMPLATE = "dn"

class SlidesState(TypedDict, total=False):
    """Estado do agente de geração de slides."""
    # Inputs
    sa_tema: str
    sa_capacidades_tecnicas: List[str]
    sa_capacidades_socioemocionais: List[str]
    sa_estrategia: str
    sa_conhecimentos: List[Dict[str, Any]]
    curso_nome: str
    uc_nome: str
    area_tecnologica: str
    num_slides: int
    template: str
    
    # Outputs
    title: str
    slides_content: List[Dict[str, Any]]
    status: str
    pptx_bytes: bytes
    error: str


def get_llm() -> AzureChatOpenAI:
    """Retorna instância do AzureChatOpenAI."""
    return AzureChatOpenAI(
        azure_deployment=os.getenv("MODEL_ID", "gpt-4o"),
        azure_endpoint=os.getenv("AZURE_OPENAI_ENDPOINT"),
        api_key=os.getenv("AZURE_OPENAI_API_KEY"),
        api_version=os.getenv("AZURE_OPENAI_API_VERSION"),
        temperature=0.7,
        max_tokens=65535,
    )


def safe_json_parse(response_text: str, fallback: Any) -> Any:
    """Tenta decodificar JSON e retorna um fallback em caso de erro."""
    text = response_text.strip()
    if text.startswith("```json"):
        text = text[7:]
    elif text.startswith("```"):
        text = text[3:]
    if text.endswith("```"):
        text = text[:-3]
    
    text = text.strip()
    try:
        start = text.find('[')
        end = text.rfind(']')
        if start != -1 and end != -1:
            text = text[start:end+1]
        return json.loads(text)
    except json.JSONDecodeError:
        logger.error(f"Erro ao decodificar JSON: {text[:100]}... Usando fallback.")
        return fallback


def _calculate_slide_distribution(total_slides: int) -> Dict[str, int]:
    """Calcula a distribuição de slides por seção."""
    total_slides = max(total_slides, 30)
    
    # Base fixa
    abertura = 2
    contextualizacao = 3
    sintese = 3
    
    base_fundamentacao = 10
    base_exemplos = 9
    base_atividades = 3
    
    # Slides extras e teto de exemplos
    MAX_EXEMPLOS = 10
    
    extras = total_slides - 30
    
    fundamentacao = base_fundamentacao
    exemplos = base_exemplos
    atividades = base_atividades
    
    if extras > 0:
        fundamentacao += round(extras * 0.60)
        exemplos += round(extras * 0.30)
        atividades += round(extras * 0.10)
        
    # Aplicar teto de exemplos
    if exemplos > MAX_EXEMPLOS:
        excedente = exemplos - MAX_EXEMPLOS
        exemplos = MAX_EXEMPLOS
        fundamentacao += excedente # Joga para fundamentação
        
    current_total = abertura + contextualizacao + fundamentacao + exemplos + atividades + sintese
    diff = total_slides - current_total
    
    # Ajuste fino
    if diff != 0:
        fundamentacao += diff
    
    return {
        "abertura": abertura,
        "contextualizacao": contextualizacao,
        "fundamentacao": fundamentacao,
        "exemplos": exemplos,
        "atividades": atividades,
        "sintese": sintese,
        "total": total_slides
    }


async def _generate_slides_section(llm, section_name: str, num_slides: int, context:  Dict[str, str], instructions: str) -> List[Dict]:
    """Gera uma seção específica de slides."""
    logger.info(f"Gerando seção: {section_name} ({num_slides} slides)")
    
    section_prompt = f"""Você é um ESPECIALISTA do SENAI.
CONTEXTO DA APRESENTAÇÃO:
{context['base_info']}

SEÇÃO ATUAL: {section_name}

⚠️ IMPORTANTE: VOCÊ DEVE GERAR EXATAMENTE {num_slides} SLIDES PARA ESTA SEÇÃO! ⚠️
NÃO GERE MENOS QUE {num_slides} SLIDES!

INSTRUÇÕES ESPECÍFICAS:
{instructions}

REGRAS DE CONTEÚDO:
1. 4 a 6 bullet points por slide (max 15 palavras cada)
2. Speaker Notes detalhadas (2-4 parágrafos) obrigatórias para TODOS os slides
3. Linguagem técnica e didática
4. QUANTIDADE EXATA: {num_slides} slides

RESPONDA APENAS COM JSON VÁLIDO (lista com EXATAMENTE {num_slides} objetos):
[
  {{"type": "...", "title": "...", "bullet_points": [...], "speaker_notes": "..."}},
  ... (total de {num_slides} objetos)
]"""

    try:
        response = await llm.ainvoke(section_prompt)
        # Parse logic
        slides = safe_json_parse(response.content, [])
        if not isinstance(slides, list):
            slides = []
            
        valid_slides = []
        for s in slides:
            if isinstance(s, dict) and "title" in s and "bullet_points" in s:
                if "speaker_notes" not in s:
                    s["speaker_notes"] = f"Explique o slide '{s.get('title')}' detalhadamente."
                valid_slides.append(s)
        
        logger.info(f"Seção {section_name}: Recebidos {len(valid_slides)} slides válidos de {num_slides} solicitados")
        
        # Completar faltantes
        slides_faltando = num_slides - len(valid_slides)
        if slides_faltando > 0:
            logger.warning(f"Seção {section_name}: Completando {slides_faltando} slides faltantes")
            for i in range(slides_faltando):
                valid_slides.append({
                    "type": "content", 
                    "title": f"{section_name} - Parte {len(valid_slides) + 1}", 
                    "bullet_points": ["Conteúdo adicional", "Aplicações práticas"],
                    "speaker_notes": "Aprofunde o debate neste ponto."
                })
            
        return valid_slides[:num_slides]
        
    except Exception as e:
        logger.error(f"Erro ao gerar seção {section_name}: {e}")
        # Fallback
        return [{
            "type": "content",
            "title": f"Seção {section_name}",
            "bullet_points": ["Conteúdo a ser desenvolvido"],
            "speaker_notes": "Erro na geração."
        }] * num_slides


async def generate_all_slides(state: SlidesState) -> Dict[str, Any]:
    """Gera todo o conteúdo dos slides com estrutura completa."""
    logger.info("Gerando conteúdo educacional dos slides...")
    llm = get_llm()
    
    num_slides = state.get("num_slides", 30)
    dist = _calculate_slide_distribution(num_slides)
    
    # Preparar contexto
    base_info = f"""TEMA: {state.get('sa_tema', 'N/A')}
CURSO: {state.get('curso_nome', 'N/A')}
UC: {state.get('uc_nome', 'N/A')}
ESTRATÉGIA: {state.get('sa_estrategia', 'Situação-Problema')}"""

    context = {"base_info": base_info}
    all_slides = []

    # 1. Capa e Título
    all_slides.append({
        "type": "cover", 
        "title": state.get('sa_tema'), 
        "subtitle": state.get('curso_nome'), 
        "speaker_notes": "Capa da apresentação."
    })
    
    all_slides.append({
        "type": "title_slide",
        "title": state.get('sa_tema'),
        "subtitle": f"{state.get('uc_nome')} - {state.get('curso_nome')}",
        "speaker_notes": "Apresentação da Situação de Aprendizagem."
    })

    # 2. Abertura
    all_slides.append({"type": "section_header", "title": "ABERTURA", "speaker_notes": "Início da seção."})
    slides_p1 = await _generate_slides_section(llm, "ABERTURA", dist['abertura'] + dist['contextualizacao'], context, 
        "Gere slides de Objetivos e Contextualização do Desafio.")
    all_slides.extend(slides_p1)

    # 3. Fundamentação
    all_slides.append({"type": "section_header", "title": "FUNDAMENTAÇÃO TEÓRICA", "speaker_notes": "Conceitos técnicos."})
    slides_p2 = await _generate_slides_section(llm, "FUNDAMENTAÇÃO TEÓRICA", dist['fundamentacao'], context, 
        "Gere slides técnicos explicativos, progressão didática.")
    all_slides.extend(slides_p2)

    # 4. Prática
    all_slides.append({"type": "section_header", "title": "APLICAÇÃO PRÁTICA", "speaker_notes": "Exemplos do mundo real."})
    slides_p3 = await _generate_slides_section(llm, "EXEMPLOS PRÁTICOS", dist['exemplos'], context, 
        "Gere exemplos de aplicação real no mercado de trabalho.")
    all_slides.extend(slides_p3)

    # 5. Atividades/Síntese
    all_slides.append({"type": "section_header", "title": "ATIVIDADES E CONCLUSÃO", "speaker_notes": "Fixação e encerramento."})
    slides_p4 = await _generate_slides_section(llm, "ATIVIDADES E SÍNTESE", dist['atividades'] + dist['sintese'], context, 
        "Gere exercícios de fixação e slides de conclusão/referências.")
    all_slides.extend(slides_p4)

    # 6. Encerramento
    all_slides.append({"type": "closing", "title": "", "bullet_points": [], "speaker_notes": "Fim."})

    return {
        "title": all_slides[0].get("title", "Apresentação"),
        "slides_content": all_slides,
        "status": "slides_generated"
    }


async def export_pptx(state: SlidesState) -> Dict[str, Any]:
    """Exporta para PPTX e preenche slides."""
    logger.info("Exportando PPTX...")
    try:
        template_key = state.get("template", DEFAULT_TEMPLATE)
        if template_key not in TEMPLATES: template_key = DEFAULT_TEMPLATE
        
        config = TEMPLATES[template_key]
        if not config["path"].exists():
             raise FileNotFoundError(f"Template não encontrado: {config['path']}")
             
        prs = Presentation(str(config["path"]))
        slides = state.get("slides_content", [])
        
        # Limpar slides existentes
        while len(prs.slides) > 0:
            rId = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[0]
            
        for idx, data in enumerate(slides):
            stype = data.get("type", "content")
            
            # Escolher layout
            layout_idx = config["layout_conteudo"]
            
            if stype == "cover": layout_idx = config["layout_capa"]
            elif stype == "title_slide": layout_idx = config.get("layout_titulo", 1)
            elif stype == "section_header": layout_idx = config.get("layout_capa_secao", 2)
            elif stype == "closing": layout_idx = config["layout_fim"]
            
            if layout_idx >= len(prs.slide_layouts): layout_idx = 0
            
            slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
            
            # Preencher
            if stype == "cover": _fill_cover_slide(slide, data, config)
            elif stype in ["title_slide", "section_header"]: _fill_section_header_slide(slide, data, config)
            elif stype == "closing": pass
            else: _fill_content_slide(slide, data, config)
            
        output = BytesIO()
        prs.save(output)
        output.seek(0)
        
        return {"pptx_bytes": output.getvalue(), "status": "completed"}
        
    except Exception as e:
        logger.error(f"Erro export PPTX: {e}", exc_info=True)
        return {"error": str(e), "status": "failed"}


def _fill_cover_slide(slide, data: Dict, config: Dict):
    # Capa visual geralmente não tem texto, mas tentamos se houver placeholder
    ph_idx = config.get("placeholder_capa", 0)
    # DN Capa (0) é imagem pura, SP Capa (0) também
    # Se for DN e layout 0, não faz nada
    if config["name"] == "Departamento Nacional" and config["layout_capa"] == 0:
        pass
    else:
        # Tenta preencher
        for shape in slide.shapes:
            if shape.has_text_frame and shape.placeholder_format.idx == ph_idx:
                shape.text_frame.text = data.get("title", "")

def _fill_section_header_slide(slide, data: Dict, config: Dict):
    title = data.get("title", "")
    subtitle = data.get("subtitle", "")
    
    # Obter placeholder específico para seção, se configurado
    section_title_idx = config.get("placeholder_section_title")
    
    # Tenta achar placeholders de título
    title_filled = False
    for shape in slide.shapes:
        if not shape.has_text_frame: continue
        ph = shape.placeholder_format
        if not ph: continue
        
        # Verifica pelo idx específico de seção ou pelo tipo TITLE (1) ou idx 0
        if (section_title_idx and ph.idx == section_title_idx) or ph.type == 1 or ph.idx == 0:
            if not title_filled:
                shape.text_frame.text = title
                title_filled = True
        elif ph.type == 2 and subtitle: # Subtitle
             shape.text_frame.text = subtitle

    if data.get("speaker_notes"):
        try: slide.notes_slide.notes_text_frame.text = data.get("speaker_notes")
        except: pass

def _fill_content_slide(slide, data: Dict, config: Dict):
    title = data.get("title", "")
    bullets = data.get("bullet_points", [])
    
    ph_title = config.get("placeholder_title", [0, 1])
    ph_content = config.get("placeholder_content", 10)
    
    # Título
    t_done = False
    for shape in slide.shapes:
        if shape.has_text_frame and shape.placeholder_format.idx in ph_title and not t_done:
            shape.text_frame.text = title
            t_done = True
            
    # Conteúdo
    for shape in slide.shapes:
        if shape.has_text_frame and shape.placeholder_format.idx == ph_content:
            tf = shape.text_frame
            tf.clear()
            for i, b in enumerate(bullets):
                p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
                p.text = str(b)
                p.level = 0
            break

    if data.get("speaker_notes"):
        try: slide.notes_slide.notes_text_frame.text = data.get("speaker_notes")
        except: pass

def create_slides_graph():
    workflow = StateGraph(SlidesState)
    workflow.add_node("generate_slides", generate_all_slides)
    workflow.add_node("export_pptx", export_pptx)
    workflow.set_entry_point("generate_slides")
    workflow.add_edge("generate_slides", "export_pptx")
    workflow.add_edge("export_pptx", END)
    return workflow.compile()

slides_agent = create_slides_graph()

async def generate_slides(
    sa_tema: str,
    sa_capacidades_tecnicas: List[str],
    sa_capacidades_socioemocionais: List[str],
    sa_estrategia: str,
    sa_conhecimentos: List[Dict[str, Any]],
    curso_nome: str,
    uc_nome: str,
    area_tecnologica: str = "",
    num_slides: int = 30,
    template: str = "dn"
) -> Dict[str, Any]:
    """Entry point."""
    if template not in TEMPLATES: template = DEFAULT_TEMPLATE
    
    initial_state: SlidesState = {
        "sa_tema": sa_tema,
        "sa_capacidades_tecnicas": sa_capacidades_tecnicas,
        "sa_capacidades_socioemocionais": sa_capacidades_socioemocionais,
        "sa_estrategia": sa_estrategia,
        "sa_conhecimentos": sa_conhecimentos,
        "curso_nome": curso_nome,
        "uc_nome": uc_nome,
        "area_tecnologica": area_tecnologica,
        "num_slides": num_slides,
        "template": template,
        "status": "starting"
    }
    
    try:
        result = await slides_agent.ainvoke(initial_state, {"recursion_limit": 50})
        return {
            "pptx_bytes": result.get("pptx_bytes"),
            "title": result.get("title", sa_tema),
            "num_slides": len(result.get("slides_content", [])),
            "status": "completed",
            "error": None
        }
    except Exception as e:
        logger.error(f"Erro slides: {e}", exc_info=True)
        return {"pptx_bytes": None, "title": None, "num_slides": 0, "status": "failed", "error": str(e)}
